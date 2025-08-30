
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt

# --- 2枚: 縦に並べて右にキャプション ---
def render_image2_default(self, slide, images, font_size):
    slide_w, slide_h = self.prs.slide_width, self.prs.slide_height
    margin = Pt(40)
    spacing = Pt(30)
    title_height = Pt(60)  # タイトル領域を固定値で確保

    # 画像の最大表示高さ
    max_h = (slide_h - margin*2 - spacing - title_height) / 2
    target_w = slide_w / 2 - margin * 2

    # 画像サイズを読み取り、リサイズ結果を保存
    scaled_sizes = []
    streams = []
    for img in images:
        stream, im = self._load_image(img["url"])
        if stream:
            iw, ih = im.size
            scale = min(max_h/ih, target_w/iw)  # 高さ基準 + 横幅制限
            new_w, new_h = iw*scale, ih*scale
            scaled_sizes.append((new_w, new_h))
            streams.append(stream)

    # 上下サイズをそろえる（小さい方に合わせる）
    min_h = min(h for _, h in scaled_sizes)
    scaled_sizes = [(w*(min_h/h), min_h) for w, h in scaled_sizes]

    # 全体の高さ（上下 + spacing）
    total_h = scaled_sizes[0][1] + scaled_sizes[1][1] + spacing
    top_start = (slide_h - total_h) / 2 + title_height/2

    shapes = []
    for i, ((new_w, new_h), stream, img) in enumerate(zip(scaled_sizes, streams, images)):
        top = top_start + i*(new_h + spacing)
        left = (slide_w/2 - new_w) / 2

        # 画像配置
        if stream:
            pic = slide.shapes.add_picture(stream, left, top, width=int(new_w), height=int(new_h))  

        # キャプションを画像の右に配置
        cap_box = slide.shapes.add_textbox(
            left+new_w+Pt(20), top, slide_w/2 - new_w - Pt(40), new_h
        )
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        self._style_text(p, img.get("caption",""), font_size, self.colors["text"])
        p.alignment = PP_ALIGN.LEFT
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        shapes.append((pic, cap_box))

    return shapes
