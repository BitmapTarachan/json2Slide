# -*- coding: utf-8 -*-
"""
majin-style slide factory with python-pptx
- AIは「設計図(JSON)」のみ生成
- この工場がテンプレートに流し込み、PPTXを安定生成
- 主要スライド型: title / section / content / two_column / compare / quote /
                  process / timeline / image / table / key_takeaways
"""
import json
import sys
import requests
import io
import os
import platform

from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement

from PIL import Image




# -------- ユーザー環境に合わせて調整可能な既定値 --------
DEFAULT_FONT = "Biz UDゴシック"           # 日本語フォントを既定化
TITLE_FONT_SIZE = Pt(36)
SUBTITLE_FONT_SIZE = Pt(18)
HEADING_FONT_SIZE = Pt(28)
BODY_FONT_SIZE = Pt(18)
CAPTION_FONT_SIZE = Pt(14)
ACA_BASE_URL = "https://myaca.azurecontainerapps.io/"

# レイアウトの既定マップ（テンプレートにより差異あり）
# 必要に応じて調整してください。
DEFAULT_LAYOUT_MAP = {
    "title": 0,           # Title Slide
    "content": 1,         # Title and Content
    "section": 2,         # Section Header
    "two_content": 3,     # Two Content
    "comparison": 4,      # Comparison
    "title_only": 5,      # Title Only
    "blank": 6,           # Blank
    "with_caption": 7,    # Content with Caption
    "pic_with_caption": 8 # Picture with Caption
}

# ---------------- Google風 CONFIG ----------------
CONFIG = {
    "BASE_PX": {"W": 960, "H": 540},

    # ---------------- テーマカラー定義 ----------------
    "COLORS": {
        "primary"    : RGBColor(0x42, 0x85, 0xF4),   # メインカラー（タイトル・強調）
        "accent"     : RGBColor(0xFB, 0xBC, 0x04),   # アクセントカラー（ハイライト・バー）
        "background" : RGBColor(0xFF, 0xFF, 0xFF),   # スライド背景
        "surface"    : RGBColor(0xF8, 0xF9, 0xFA),   # セクション背景・ボックス背景
        "text"       : RGBColor(0x33, 0x33, 0x33),   # 標準本文
        "subtext"    : RGBColor(0x9E, 0x9E, 0x9E),   # 補助テキスト（日付・キャプション）
        "ghost"      : RGBColor(0xEF, 0xEF, 0xED),   # ゴースト数字・区切り用
    },

    # ---------------- フォント定義 ----------------
    "FONTS": {
        "family": "Biz UDゴシック",
        "sizes": {
            "title"        : Pt(45),
            "sectionTitle" : Pt(38),
            "contentTitle" : Pt(30),
            "subhead"      : Pt(28),
            "body"         : Pt(22),
            "caption"      : Pt(18),
            "ghostNum"     : Pt(180),
        }
    },

    # ---------------- レイアウト座標 ----------------
    "POS_PX": {
        "titleSlide": {
            "subject":  { "left": 80, "top": 140, "width": 800, "height": 40},  
                "title"    : { "left": 80, "top": 190, "width": 800, "height": 90 },  
                "lecturer" : { "left": 80, "top": 290, "width": 400, "height": 40 },  
                "date"     : { "left": 80, "top": 330, "width": 250, "height": 40 },  
        },        
        "sectionSlide": {
            "title"    : { "left":  55, "top": 230, "width": 840, "height":  80 },
            "ghostNum" : { "left": 100, "top": 120, "width": 300, "height": 200 },
        },
        "contentSlide": {
            "title"    : { "left": 28, "top":  30, "width": 830, "height":  50 },
            "subhead"  : { "left": 25, "top": 100, "width": 830, "height":  30 },
            "body"     : { "left": 25, "top": 150, "width": 910, "height": 303 },
        }
    }
}

THEME_COLORS = {
    "Default": {
        "primary"    : RGBColor(0x42, 0x85, 0xF4),  # Googleブルー
        "accent"     : RGBColor(0xFB, 0xBC, 0x04),  # 黄色
        "background" : RGBColor(0xFF, 0xFF, 0xFF),
        "surface"    : RGBColor(0xF8, 0xF9, 0xFA),
        "text"       : RGBColor(0x33, 0x33, 0x33),
        "subtext"    : RGBColor(0x9E, 0x9E, 0x9E),
        "ghost"      : RGBColor(0xEF, 0xEF, 0xED),
    },
    "Nature": {
        "primary"    : RGBColor(0x2E, 0x7D, 0x32),  # 深緑
        "accent"     : RGBColor(0xFF, 0xA0, 0x00),  # オレンジ
        "background" : RGBColor(0xFF, 0xFF, 0xF5),
        "surface"    : RGBColor(0xE8, 0xF5, 0xE9),
        "text"       : RGBColor(0x1B, 0x5E, 0x20),
        "subtext"    : RGBColor(0x6D, 0x6D, 0x6D),
        "ghost"      : RGBColor(0xC8, 0xE6, 0xC9),
    },
    "Dark": {
        "primary"    : RGBColor(0xBB, 0x86, 0xFC),  # 紫
        "accent"     : RGBColor(0x03, 0xDA, 0xC6),  # シアン
        "background" : RGBColor(0x12, 0x12, 0x12),
        "surface"    : RGBColor(0x1E, 0x1E, 0x1E),
        "text"       : RGBColor(0xEE, 0xEE, 0xEE),
        "subtext"    : RGBColor(0xAA, 0xAA, 0xAA),
        "ghost"      : RGBColor(0x33, 0x33, 0x33),
    },
    "Monochrome": {
        "primary"    : RGBColor(0x00, 0x00, 0x00),
        "accent"     : RGBColor(0x55, 0x55, 0x55),
        "background" : RGBColor(0xFF, 0xFF, 0xFF),
        "surface"    : RGBColor(0xF0, 0xF0, 0xF0),
        "text"       : RGBColor(0x00, 0x00, 0x00),
        "subtext"    : RGBColor(0x77, 0x77, 0x77),
        "ghost"      : RGBColor(0xDD, 0xDD, 0xDD),
    }
}

# ---------------- Layout Manager ----------------
class LayoutManager:
    def __init__(self, config):
        self.cfg = config
        self.base_w = config["BASE_PX"]["W"] * 0.75
        self.base_h = config["BASE_PX"]["H"] * 0.75
        self.page_w = Inches(13.33)  # 16:9 幅
        self.page_h = Inches(7.5)    # 16:9 高さ
        self.scale_x = self.page_w / self.base_w
        self.scale_y = self.page_h / self.base_h

    def get_rect(self, path: str):
        keys = path.split(".")
        pos = self.cfg["POS_PX"]
        for k in keys:
            pos = pos[k]
        def px2pt(px): return px * 0.75
        return {
            "left": px2pt(pos["left"]) * self.scale_x,
            "top": px2pt(pos["top"]) * self.scale_y,
            "width": px2pt(pos["width"]) * self.scale_x,
            "height": px2pt(pos["height"]) * self.scale_y,
        }

# ---------------------- ユーティリティ ----------------------

def set_paragraph_style(paragraph, text: str, font_size: Pt, bold=False, italic=False, color: Optional[RGBColor]=None, align=None):
    paragraph.text = text
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.name = DEFAULT_FONT
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if align:
        paragraph.alignment = align


def set_text_frame_bullets(tf, lines: List[str], level: int = 0):
    """最初の段落を上書きし、以降は追加。"""
    tf.clear()
    if not lines:
        return
    p0 = tf.paragraphs[0]
    set_paragraph_style(p0, lines[0], BODY_FONT_SIZE)
    p0.level = level
    for line in lines[1:]:
        p = tf.add_paragraph()
        set_paragraph_style(p, line, BODY_FONT_SIZE)
        p.level = level

def add_speaker_notes(slide, notes: Optional[str]):
    if not notes:
        return
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    set_paragraph_style(p, notes, Pt(14))

def hex_to_rgbcolor(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)

def ensure_list(x) -> List[Any]:
    if x is None:
        return []
    return x if isinstance(x, list) else [x]


# ---------------- Slide Factory ----------------
class SlideFactory:
    def __init__(self, plan, theme, config=CONFIG):

        self.plan = plan
        self.theme = theme
        self.prs = Presentation()

        self.image_base_dir = "images"
        self.is_aca = False
        if platform.system() == "Windows":
            self.is_aca = False
        else:
            self.is_aca = True
        
        # イメージキャッシュ
        self._image_cache = {}

        # カラーテーマ選択
        theme_name = plan.get("color-theme", "Default")

        if theme_name == "Custom" and "colors" in plan:
            self.colors = {
                k: hex_to_rgbcolor(v) for k, v in plan["colors"].items()
            }
        else:
            self.colors = THEME_COLORS.get(theme_name, THEME_COLORS["Default"])

        self.config = config
        self.fonts = config["FONTS"]
        self.layout = LayoutManager(config)

        # 16:9 に固定
        self.prs.slide_width = self.layout.page_w
        self.prs.slide_height = self.layout.page_h
        
        # 全体背景
        self._global_bg = None
        if plan.get("background-image"):
            self._global_bg, _ = self._load_image(plan["background-image"])


    def add_slide(self,spec):
        t = spec.get("type")
        if t == "title":
            return self.theme.render_title(self, spec) 
        elif t == "section":
            return self.theme.render_section(self, spec) 
        elif t == "content":
            return self.theme.render_content(self, spec)
        elif t == "cards" :
            return self.theme.render_cards(self, spec) 
        elif t == "compare":
            return self.theme.render_compare(self, spec)
        elif t == "progress":
            return self.theme.render_progress(self, spec)
        elif t == "timeline":
            return self.theme.render_timeline(self, spec)
        elif t == "image-auto":
            return self.theme.render_image_auto(self, spec)
        elif t == "qa-question":
            return self.theme.render_qa_question(self, spec)
        elif t == "qa-answer":
            return self.theme.render_qa_answer(self, spec)
        elif t == "table":
            return self.theme.render_table(self, spec)
        elif t == "flow":
            return self.theme.render_flow(self, spec)
        elif t == "highlight":
            return self.theme.render_highlight(self, spec)
        elif t == "hero":
            return self.theme.render_hero(self, spec) 
        elif t == "quote":
            return self.theme.render_quote(self, spec)


    def save(self, out_path: str):
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(out_path)

    # ---------------- 内部ユーティリティ ----------------
    def _new_slide(self, data: Dict[str, Any],apply_background = True):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # 背景
        self._apply_background(s,data,apply_background)
        # スライドノート
        note_text = data.get("note", "")
        s.notes_slide.notes_text_frame.text = note_text
        return s

    def _style_text(self, paragraph, text: str, size: Pt, bold=False,
                    italic=False, color=None, align=None):
        """統一的にフォントスタイルを適用"""
        paragraph.text = text
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.name = self.fonts["family"]
        run.font.size = size
        run.font.bold = bool(bold) if bold is not None else None
        run.font.italic = italic
        run.font.color.rgb = color or self.colors["text"]
        if align:
            paragraph.alignment = align

    def _load_image(self, path_or_url: str):

        # ACAかローカルのパスを解釈
        if not path_or_url.startswith(("http://","https://")): 
            # ファイル名らしきものの解釈
            if self.is_aca:
                path_or_url = f"{ACA_BASE_URL}/{self.image_base_dir}/{path_or_url}" 
            else:
                path_or_url = os.path.join(self.image_base_dir, path_or_url)
        
        try:
            # キャッシュヒット
            if path_or_url in self._image_cache:
                # BytesIOは再利用のため毎回seek(0)して返す
                stream, im = self._image_cache[path_or_url]
                stream.seek(0)
                return stream, im

            # URLから取得
            if path_or_url.startswith(("http://", "https://")):
                response = requests.get(path_or_url, timeout=10)
                response.raise_for_status()
                stream = io.BytesIO(response.content)
            else:
                # ローカルファイル
                with open(path_or_url, "rb") as f:
                    stream = io.BytesIO(f.read())

            # PILでロードしてキャッシュ
            im = Image.open(stream)
            im.load()           # Lazy読み込みを強制
            stream.seek(0)      # 再利用に備えて戻す
            self._image_cache[path_or_url] = (stream, im)

            return stream, im
        except Exception as e:
            print(f"[WARN] 画像を読み込めませんでした: {path_or_url} ({e})")
            return None
    
    def _add_slide_title(self, slide, title: str):
        """
        スライドタイトルを描画する共通関数
        左端にアクセントカラーの縦長バーを置き、
        その右にタイトルテキストを配置する。
        """
        # レイアウトからタイトル領域を取得
        t_rect = self.layout.get_rect("contentSlide.title")
        left, top, width, height = t_rect["left"], t_rect["top"], t_rect["width"], t_rect["height"]

        # --- 縦長バー ---
        bar_width = Pt(6)   # 適度に細いバー
        bar_margin = Pt(4)  # テキストとの間隔

        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top,
            bar_width, height
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = self.colors["accent"]
        slide.shapes[-1].line.fill.background()  # 枠線なし
        slide.shapes[-1].shadow.inherit = False  # 影を消す（フラット）

        # --- タイトルテキスト ---
        tbox = slide.shapes.add_textbox(
            left + bar_width + bar_margin, top,
            width - bar_width - bar_margin, height
        )
        tf = tbox.text_frame
        tp = tf.paragraphs[0]
        self._style_text(
            tp,
            title,
            self.fonts["sizes"]["contentTitle"],
            bold=True,
            color=self.colors["primary"]
        )
        tp.alignment = PP_ALIGN.LEFT
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _apply_background(self, slide, slide_data: dict = None,apply_background = True):
        bg_url = None
        if slide_data and "background-image" in slide_data:
            bg_url = slide_data["background-image"]
            stream, _ = self._load_image(bg_url)
            if stream :
                slide.shapes.add_picture(stream, 0, 0,
                                        width=self.prs.slide_width,
                                        height=self.prs.slide_height)
        elif apply_background and self._global_bg:
            if self._global_bg :
                slide.shapes.add_picture(self._global_bg, 0, 0,
                                        width=self.prs.slide_width,
                                        height=self.prs.slide_height)
        else:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self.colors["background"]    
            
    def _set_shape_transparency(self, shape, alpha_val: int):
        """
        shape.fill に alpha 要素を追加して透過を指定（alpha_val は 0〜100000）
        例: 40000 = 40% 透過
        """
        from pptx.oxml.xmlchemy import OxmlElement
        ts = shape.fill._xPr.solidFill
        sF = ts.get_or_change_to_srgbClr()
        alpha_elem = OxmlElement('a:alpha')
        alpha_elem.set('val', str(alpha_val))
        sF.append(alpha_elem)

    # ---------------------- ビルド関数 ----------------------
def build_pptx_from_plan(plan: Dict[str, Any], out_path: str):
    
    from themes_default import DefaultTheme
    from themes_simplenote import SimpleNoteTheme

    theme = SimpleNoteTheme()
    sf = SlideFactory(plan,theme)

    for spec in plan.get("slides", []):
        sf.add_slide(spec)

    sf.save(out_path)

# ---------------- CLI ----------------
if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python json2.py plan.json out.pptx")
        sys.exit(1)

    plan_path = Path(sys.argv[1])
    out_path = sys.argv[2]
    with plan_path.open("r", encoding="utf-8") as f:
        plan = json.load(f)

    build_pptx_from_plan(plan, out_path)
    print(f"✅ Done: {out_path}")
