from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt

# --- 4枚: 横4グリッド + 下キャプション ---
def render_image4_default(self, slide, images, font_size):
    slide_w, slide_h = self.prs.slide_width, self.prs.slide_height
    margin = Pt(40)
    spacing = Pt(30)
    n = 4

    # 横方向の基準幅（スライド幅の90%に収める）
    max_total_w = slide_w * 0.9
    target_w = (max_total_w - spacing * (n - 1)) / n

    resized = []
    streams = []

    for img in images:
        stream, im = self._load_image(img["url"])
        if stream:
            iw, ih = im.size
            scale = target_w / iw
            new_w, new_h = iw * scale, ih * scale
            resized.append((new_w, new_h, img))
            streams.append(stream)

    # 最大高さを取得（キャプション基準にする）
    max_h = max(h for _, h, _ in resized)

    # 横方向の開始位置（中央寄せ）
    total_w = sum(w for w, _, _ in resized) + spacing * (n - 1)
    left_base = (slide_w - total_w) / 2
    top_img = slide_h * 0.35

    shapes = []
    cur_left = left_base
    for (w, h, img), stream in zip(resized, streams):
        # 画像（下揃え）
        pic_top = top_img + (max_h - h)
        if stream: 
            pic = slide.shapes.add_picture(stream, cur_left, pic_top, width=int(w), height=int(h))

        # キャプション（全画像上揃え）
        caption = img.get("caption", "")
        cap_box = slide.shapes.add_textbox(cur_left, top_img + max_h + Pt(10), w, Pt(40))
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        self._style_text(p, caption, font_size, color=self.colors["text"])
        p.alignment = PP_ALIGN.LEFT

        if stream:
            shapes.append((pic, cap_box))

        cur_left += w + spacing

    return shapes
