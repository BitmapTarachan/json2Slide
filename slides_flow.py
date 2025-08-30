from typing import Any, Dict

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 手順解説
def render_flow_default(factory, data: Dict[str, Any]):
    s = factory._new_slide(data)
    factory._add_slide_title(s, data["title"])

    steps = data.get("steps", [])
    body = data.get("bodyText", "")
    n = len(steps)
    direction = data.get("direction", "horizontal")

    slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height
    margin = Pt(60)
    spacing = Pt(50)

    if direction == "horizontal":
        # 横フロー
        box_w = (slide_w - margin*2 - spacing*(n-1)) / n
        box_h = Pt(120)
        top = slide_h/2 - box_h/2
        left = margin

        for i, text in enumerate(steps):
            # ラウンドボックス
            shape = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_w, box_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = factory.colors["surface"]
            shape.line.color.rgb = factory.colors["primary"]

            # 数字（左上に重ねる）
            num_box = s.shapes.add_textbox(left-Pt(20), top-Pt(40), Pt(40), Pt(40))
            tf_num = num_box.text_frame
            tf_num.text = str(i+1)
            p_num = tf_num.paragraphs[0]
            run_num = p_num.runs[0]
            run_num.font.size = Pt(55)
            run_num.font.bold = True
            run_num.font.color.rgb = factory.colors["accent"]
            p_num.alignment = PP_ALIGN.LEFT

            # 本文
            tf = shape.text_frame
            tf.text = text
            p = tf.paragraphs[0]
            run = p.runs[0]
            run.font.size = Pt(20)
            run.font.name = "BIZ UDPゴシック"
            run.font.color.rgb = factory.colors["text"]  
            p.alignment = PP_ALIGN.CENTER

            # 矢印
            if i < n-1:
                arrow = s.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    left+box_w+5, top+box_h/3, spacing-10, Pt(40)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = factory.colors["accent"]
                arrow.line.fill.background()

            left += box_w + spacing

        # BodyText（下）
        if body:
            tbox = s.shapes.add_textbox(
                Pt(60), top+box_h+spacing, slide_w-Pt(120), Pt(100)
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            factory._style_text(p, body, factory.fonts["sizes"]["body"], color=factory.colors["text"])

    else:
        # 縦フロー（左寄せ）
        flow_area_w = slide_w * 0.55   # 左2/3
        body_area_left = slide_w * 0.65
        box_w = flow_area_w * 0.9      # さらに少し狭く
        box_h = (slide_h - margin*2 - spacing*(n-1)) / n
        left = margin + Pt(40)
        top = margin + Pt(40)

        for i, text in enumerate(steps):
            # ラウンドボックス
            shape = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_w, box_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = factory.colors["surface"]
            shape.line.color.rgb = factory.colors["primary"]

            # 数字（ボックスの左外）
            num_box = s.shapes.add_textbox(left-Pt(60), top, Pt(50), box_h)
            tf_num = num_box.text_frame
            tf_num.text = str(i+1)
            p_num = tf_num.paragraphs[0]
            run_num = p_num.runs[0]
            run_num.font.size = Pt(36)
            run_num.font.bold = True
            run_num.font.color.rgb = factory.colors["accent"]
            p_num.alignment = PP_ALIGN.CENTER

            # 本文
            tf = shape.text_frame
            tf.text = text
            p = tf.paragraphs[0]
            run = p.runs[0]
            run.font.size = Pt(20)
            run.font.name = "BIZ UDPゴシック"
            run.font.color.rgb = factory.colors["text"] 
            p.alignment = PP_ALIGN.CENTER

            # 矢印（下向き）
            if i < n-1:
                arrow = s.shapes.add_shape(
                    MSO_SHAPE.DOWN_ARROW,
                    left + box_w / 2 - Pt(20), top+box_h+5, Pt(40), spacing-10
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = factory.colors["accent"]
                arrow.line.fill.background()

            top += box_h + spacing

        # BodyText（右1/3）
        if body:
            body_box = s.shapes.add_textbox(
                body_area_left, margin + Pt(40), slide_w - body_area_left - Pt(40),
                slide_h - margin*2
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            factory._style_text(p, body, factory.fonts["sizes"]["body"], color=factory.colors["text"])
