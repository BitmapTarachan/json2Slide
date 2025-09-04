from typing import List
from pptx.util import Pt
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE

def render_content_default(factory, data):
    s = factory._new_slide(data)
    factory._add_slide_title(s, data.get("title", ""))

    # Subhead
    subhead = data.get("subhead")
    if subhead:
        s_rect = factory.layout.get_rect("contentSlide.subhead")
        sbox = s.shapes.add_textbox(s_rect["left"], s_rect["top"], s_rect["width"], s_rect["height"])
        sp = sbox.text_frame.paragraphs[0]
        factory._style_text(
            sp,
            subhead,
            factory.fonts["sizes"]["subhead"],
            color=factory.colors["subtext"]
        )
        sbox.text_frame.word_wrap = True

    # Body area base rect
    b_rect = factory.layout.get_rect("contentSlide.body")

    points: List[str] = data.get("points", [])
    body_text: str = data.get("bodyText", "")

    last_y = b_rect["top"]

    # --- 箇条書き ---
    if points:
        # フォントサイズと行間を元に高さを見積もり
        line_h = factory.fonts["sizes"]["body"].pt * 1.6  # だいたい1.6倍行間
        box_h = Pt(line_h * len(points))

        bbox = s.shapes.add_textbox(b_rect["left"]+ Cm(0.5), b_rect["top"], b_rect["width"], box_h)
        tf = bbox.text_frame
        tf.clear()
        tf.word_wrap = True

        for i, line in enumerate(points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            factory._style_text(p, f"・{line}", factory.fonts["sizes"]["body"], color=factory.colors["text"],bold=True)
            p.space_after = Pt(6)

        last_y = b_rect["top"] + box_h + Cm(0.5)

    # --- 長文 ---
    if body_text:
        # 背景ボックス（薄い色）
        body_h = Cm(4.5)  # デフォルト高さ、必要に応じて調整
        bg = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            b_rect["left"],
            last_y,
            b_rect["width"],
            body_h
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = factory.colors["ghost"]  # ← 薄めの色を使う想定
        bg.line.color.rgb = factory.colors["ghost"]       # 枠線消し
        bg.shadow.inherit = False

        # 本文テキスト
        lbox = s.shapes.add_textbox(b_rect["left"] + Cm(0.5), last_y + Cm(0.5),
                                    b_rect["width"] - Cm(1.0), body_h - Cm(1.0))
        tf2 = lbox.text_frame
        tf2.clear()
        tf2.word_wrap = True
        lp = tf2.paragraphs[0]
        factory._style_text(lp, body_text, factory.fonts["sizes"]["body"], color=factory.colors["text"])

    return s