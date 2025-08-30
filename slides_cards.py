
from typing import Any, Dict

from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm
from pptx.enum.text import MSO_ANCHOR

def render_cards_default(factory, data: Dict[str, Any]):
    """カード形式スライド"""
    s = factory._new_slide(data)
    factory._add_slide_title(s, data.get("title","一覧"))

    items = data.get("items", [])
    cols = min(3, max(1, int(data.get("columns", 3))))
    gap = Cm(0.5)
    card_w = (factory.prs.slide_width - Cm(2) - gap * (cols - 1)) / cols
    rows = (len(items) + cols - 1) // cols
    card_h = Cm(5)

    for idx, item in enumerate(items):
        r, c = divmod(idx, cols)
        left = Cm(1) + c * (card_w + gap)
        top = Cm(4) + r * (card_h + gap)

        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = factory.colors["surface"]
        card.line.color.rgb = factory.colors["ghost"]

        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP  

        if isinstance(item, dict):
            title = str(item.get("title", ""))
            desc = str(item.get("desc", ""))
            p = tf.paragraphs[0]
            factory._style_text(p, title, factory.fonts["sizes"]["body"], bold=True, color=factory.colors["primary"])
            if desc:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(8) 
                factory._style_text(p2, desc, factory.fonts["sizes"]["caption"], color=factory.colors["text"])
        else:
            p = tf.paragraphs[0]
            factory._style_text(p, str(item), factory.fonts["sizes"]["body"], color=factory.colors["text"])

    return s
