

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm

def render_compare_default(factory, data):
    s = factory._new_slide(data)
    factory._add_slide_title(s, data.get("title","比較"))

    # ボックス配置
    margin = Cm(1.5)
    gap = Cm(1.5)
    box_w = (factory.prs.slide_width - margin * 2 - gap) / 2
    box_h = Cm(8)
    top = Cm(4)

    def add_box(x, title, items):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = factory.colors["surface"]
        box.line.color.rgb = factory.colors["ghost"]

        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True

        # タイトル
        p = tf.paragraphs[0]
        factory._style_text(p, title, factory.fonts["sizes"]["subhead"], bold=True, color=factory.colors["text"])
        p.space_after = Pt(15)

        # 箇条書き
        for item in items:
            para = tf.add_paragraph()
            factory._style_text(para, f"• {item}", factory.fonts["sizes"]["body"], color=factory.colors["text"])
            para.space_after = Pt(6)

    # 左ボックス
    add_box(margin,
            data.get("leftTitle", "選択肢A"),
            data.get("leftItems", [])
    )

    # 右ボックス
    add_box(margin + box_w + gap,
            data.get("rightTitle", "選択肢B"),
            data.get("rightItems", [])
    )

    # --- 結論 BodyText ---
    body_text = data.get("bodyText", "")
    if body_text:
        b_rect = factory.layout.get_rect("contentSlide.body")
        lbox = s.shapes.add_textbox(b_rect["left"], top + box_h + Cm(1.0), b_rect["width"], Cm(3))
        tf2 = lbox.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        factory._style_text(p, body_text, factory.fonts["sizes"]["body"], color=factory.colors["text"], align=PP_ALIGN.CENTER)

    return s
