#slide_section.py
from pptx.enum.text import PP_ALIGN

def render_section_default(factory, data):
    s = factory._new_slide(data)

    # ゴースト番号
    g_rect = factory.layout.get_rect("sectionSlide.ghostNum")
    gbox = s.shapes.add_textbox(g_rect["left"], g_rect["top"], g_rect["width"], g_rect["height"])
    gp = gbox.text_frame.paragraphs[0]
    factory._style_text(
        gp,
        str(data.get("sectionNo", "01")),
        factory.fonts["sizes"]["ghostNum"],
        bold=True,
        color=factory.colors["ghost"]
    )

    # セクションタイトル
    t_rect = factory.layout.get_rect("sectionSlide.title")
    tbox = s.shapes.add_textbox(t_rect["left"], t_rect["top"], t_rect["width"], t_rect["height"])
    tp = tbox.text_frame.paragraphs[0]
    factory._style_text(
        tp,
        data.get("title", ""),
        factory.fonts["sizes"]["sectionTitle"],
        bold=True,
        color=factory.colors["text"],
        align=PP_ALIGN.CENTER
    )
    return s
