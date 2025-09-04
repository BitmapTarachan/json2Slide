from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

def render_features_default(factory, data):
    s = factory._new_slide(data)
    factory._add_slide_title(s, data.get("title", "特徴"))

    items = data.get("items", [])
    n = min(len(items), 4)
    if n == 0:
        return s

    # レイアウト計算（間隔を狭め、中央寄せ）
    outer_margin = Cm(3.0)  # 左右の外側余白を多めに
    gap = Cm(1.0)           # アイテム間隔は狭め
    box_w = (factory.prs.slide_width - outer_margin * 2 - gap * (n - 1)) / n
    box_h = Cm(9)
    top = Cm(4)
    title_fontsize = Pt(22) if n <= 3 else Pt(18)
    desc_fontsize  = Pt(18) if n <= 3 else Pt(16)

    title_h = Cm(2)
    desc_h = box_h - title_h - Cm(1.0)

    for i, item in enumerate(items[:n]):
        x = outer_margin + i * (box_w + gap)

        # 背景ボックス
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = factory.colors["surface"]
        box.line.color.rgb = factory.colors["ghost"]
        box.shadow.inherit = False

        # タイトル
        tbox = s.shapes.add_textbox(x + Cm(0.5), top + Cm(0.5), box_w - Cm(1.0), title_h)
        tf1 = tbox.text_frame
        tf1.clear()
        tf1.word_wrap = True  # 折り返し有効
        p1 = tf1.paragraphs[0]
        factory._style_text(
            p1,
            item.get("title", ""),
            title_fontsize,
            bold=True,
            color=factory.colors["text"]
        )

        # 説明文
        dbox = s.shapes.add_textbox(x + Cm(0.5), top + title_h + Cm(1.0), box_w - Cm(1.0), desc_h)
        tf2 = dbox.text_frame
        tf2.clear()
        tf2.word_wrap = True  # 折り返し有効
        desc = item.get("desc", "")
        if desc:
            p2 = tf2.paragraphs[0]
            factory._style_text(
                p2,
                desc,
                desc_fontsize,
                color=factory.colors["text"]
            )
    
    note = data.get("note", "")
    if note:
        # Box群の下に横幅いっぱいで追加
        note_top = top + box_h + Cm(1.0)
        note_box = s.shapes.add_textbox(outer_margin, note_top, 
                                        factory.prs.slide_width - outer_margin * 2, Cm(3))
        tf_note = note_box.text_frame
        tf_note.clear()
        tf_note.word_wrap = True
        p_note = tf_note.paragraphs[0]
        factory._style_text(
            p_note,
            note,
            Pt(18),
            color=factory.colors["text"],
            align=PP_ALIGN.LEFT
        )

    return s
