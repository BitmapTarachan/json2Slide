

from typing import Any, Dict

from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm

from PIL import Image

# --- Progress ---
def render_progress_default(factory, data: Dict[str, Any]):
    """進捗バー"""
    s = factory._new_slide(data)

    factory._add_slide_title(s, data.get("title","進捗状況"))

    items = data.get("items", [])
    bar_left = Cm(6)
    bar_width = Cm(23) 
    bar_height = Cm(1)
    v_gap = Cm(1.5)

    for i, item in enumerate(items):
        label = str(item.get("label", f"Step {i+1}"))
        pct = max(0, min(100, int(item.get("percent", 0))))

        y = Cm(4) + i * v_gap

        # ラベル
        lbox = s.shapes.add_textbox(Cm(0.5), y - Cm(0.2), bar_left - Cm(1), bar_height)
        lp = lbox.text_frame.paragraphs[0]
        factory._style_text(lp, label, factory.fonts["sizes"]["body"], color=factory.colors["text"], align=PP_ALIGN.RIGHT)

        # 背景バー
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, bar_width, bar_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = factory.colors["ghost"]
        bg.line.fill.background()

        # 実バー
        fg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, bar_width * pct / 100, bar_height)
        fg.fill.solid()
        fg.fill.fore_color.rgb = factory.colors["accent"]
        fg.line.fill.background()

        # 数値ラベル（右側も余裕広く）
        pbox = s.shapes.add_textbox(bar_left + bar_width + Cm(0.5), y - Cm(0.2), Cm(4), bar_height)
        pp = pbox.text_frame.paragraphs[0]
        factory._style_text(pp, f"{pct}%", factory.fonts["sizes"]["body"], color=factory.colors["subtext"])

    return s
