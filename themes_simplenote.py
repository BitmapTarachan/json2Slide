# themes_simplenote.py
import themes_base
import slides_section
import slides_content
import slides_cards
import slides_compare
import slides_progress
import slides_timeline
import slides_image1
import slides_image2
import slides_image3
import slides_image4
import slides_qa_question
import slides_qa_answer
import slides_table
import slides_flow
import slides_highlight
import slides_quote
import slides_hero
import slides_features
import slides_closing

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class SimpleNoteTheme(themes_base.SlideTheme):


    def add_full_height_image(self, factory, slide):
        stream, im = factory._load_image("simplenote1.png")
        if stream is None:
            print("[WARN] 画像を読み込めません: simplenote1.png")
            return None

        # スライド全体の高さ
        slide_height = factory.prs.slide_height

        # 画像を追加
        slide.shapes.add_picture(stream, 0, 0, height=slide_height)
    
    def top_title(self, slide, factory, title_str):
        if title_str:
            sbox = slide.shapes.add_textbox(Pt(100), Pt(20), factory.prs.slide_width - Pt(100), Pt(32))
            sp = sbox.text_frame.paragraphs[0]
            factory._style_text(
                sp,
                title_str,
                Pt(22),
                color=factory.colors["text"],
                bold=True
            )
            sbox.text_frame.word_wrap = True
        
        # タイトル下の横線
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Pt(70), Pt(60), Pt(850), Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 黒
        line.line.fill.background()  # 枠線なし
        line.shadow.inherit = False
        
    def render_title(self, factory, data):
        slide = factory._new_slide(data)
        
        # 横棒
        self.add_full_height_image(factory,slide)

        #subject
        subject = data.get("subject")
        if subject:
            sbox = slide.shapes.add_textbox(Pt(100), Pt(170), factory.prs.slide_width - Pt(100), Pt(30))
            sp = sbox.text_frame.paragraphs[0]
            factory._style_text(
                sp,
                subject,
                Pt(24),
                color=factory.colors["text"]
            )
            sbox.text_frame.word_wrap = True
            

        # タイトル
        head = data.get("title")
        if head:
            sbox = slide.shapes.add_textbox(Pt(100), Pt(200), factory.prs.slide_width - Pt(100), Pt(50))
            sp = sbox.text_frame.paragraphs[0]
            factory._style_text(
                sp,
                head,
                factory.fonts["sizes"]["contentTitle"],
                color=factory.colors["text"],
                bold=True
            )
            sbox.text_frame.word_wrap = True
        # タイトル下の横線
        line_top = Pt(260)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Pt(100), line_top, Pt(300), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 黒
        line.line.fill.background()  # 枠線なし
        line.shadow.inherit = False

        # タイトル下の横線2
        line_top = Pt(260)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Pt(100), line_top + Pt(1), Pt(850), Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 黒
        line.line.fill.background()  # 枠線なし
        line.shadow.inherit = False

        # 講師名（固定）
        l_rect = factory.layout.get_rect("titleSlide.lecturer")
        lbox = slide.shapes.add_textbox(Pt(100), l_rect["top"], l_rect["width"], l_rect["height"])
        lp = lbox.text_frame.paragraphs[0]
        factory._style_text(
            lp,
            "講師名：〇〇　〇〇",
            factory.fonts["sizes"]["body"],
            color=factory.colors["text"]
        )
        # 日付
        d_rect = factory.layout.get_rect("titleSlide.date")
        dbox = slide.shapes.add_textbox(Pt(100), d_rect["top"], d_rect["width"], d_rect["height"])
        dp = dbox.text_frame.paragraphs[0]
        factory._style_text(
            dp,
            data.get("date", ""),
            factory.fonts["sizes"]["body"],
            color=factory.colors["subtext"]
        )

    def render_section(self, factory, data):
        return slides_section.render_section_default(factory, data)
    
    def render_content(self, factory, data):
        slide = slides_content.render_content_default(factory, data)
        self.delete_default_title(slide)

        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)
        # 横棒
        self.add_full_height_image(factory,slide)

        return slide
    
    def render_cards(self, factory, data):
        slide = slides_cards.render_cards_default(factory, data)
        self.delete_default_title(slide)

        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)
    
    def render_compare(self, factory, data):
        slide = slides_compare.render_compare_default(factory, data)
        self.delete_default_title(slide)
        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)
        # 横棒
        self.add_full_height_image(factory,slide)

        return slide

    
    def render_progress(self, factory, data):
        slide = slides_progress.render_progress_default(factory, data)
        self.delete_default_title(slide)
        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)
   
    def render_timeline(self, factory, data):
        slide = slides_timeline.render_timeline_default(factory, data)
        self.delete_default_title(slide)
        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)

        return slide

    
    def render_image1(self, factory, slide, images, font_size):
        return slides_image1.render_image1_default(factory, slide, images, font_size)

    def render_image2(self, factory, slide, images, font_size):
        return slides_image2.render_image2_default(factory, slide, images, font_size)

    def render_image3(self, factory, slide, images, font_size):
        return slides_image3.render_image3_default(factory, slide, images, font_size)

    def render_image4(self, factory, slide, images, font_size):
        return slides_image4.render_image4_default(factory, slide, images, font_size)

    def render_qa_question(self, factory, data):
        slides_qa_question.render_qa_question_defaults(factory, data)
    
    def render_qa_answer(self, factory, data):
        return slides_qa_answer.render_qa_answer_default(factory, data)

    def render_table(self, factory, data):
        slide = factory._new_slide(data)
        slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height

        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)

        # 横棒
        self.add_full_height_image(factory,slide)

        headers = data.get("headers", [])
        rows = data.get("rows", [])
        n_rows, n_cols = len(rows) + 1, len(headers)        

        top = Pt(100)
        left = Pt(100)
        width = int(slide_w - Pt(160))
        height = int(slide_h * 0.55)

        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table

        # 列幅（整数化必須）
        col_width = int(width / n_cols)
        for col in table.columns:
            col.width = col_width

        # 行高さ（整数化必須）
        row_height = int(height / n_rows)
        for row in table.rows:
            row.height = row_height

        # ヘッダー
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = header
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.name = "BIZ UDゴシック"
            run.font.color.rgb = factory.colors["background"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(128,128,128)
            p.alignment = PP_ALIGN.CENTER

        # データ
        for i, row_data in enumerate(rows):
            for j, val in enumerate(row_data):
                cell = table.cell(i + 1, j)
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
                run.font.size = Pt(16)
                run.font.name = "BIZ UDゴシック"
                run.font.color.rgb = factory.colors["text"]
                p.alignment = PP_ALIGN.CENTER
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = factory.colors["surface"]
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = factory.colors["background"]

        # bodyText（高さが溢れないように制限）
        body_text = data.get("bodyText")
        if body_text:
            b_top = min(top + height + Pt(20), slide_h - Pt(100))
            b_left = Pt(100)
            b_width = int(slide_w - Pt(150))
            b_height = int(slide_h - b_top - Pt(40))

            box = slide.shapes.add_textbox(b_left, b_top, b_width, b_height)
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = body_text
            run.font.size = Pt(16)
            run.font.name = "BIZ UDゴシック"
            run.font.color.rgb = factory.colors["text"]

        return slide
    
    def render_flow(self, factory, data):
        slide = slides_flow.render_flow_default(factory, data)
        self.delete_default_title(slide)
        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)
        # 横棒
        self.add_full_height_image(factory,slide)

        return slide


    def render_highlight(self, factory, data):
        slide = slides_highlight.render_highlight_default(factory, data)
        self.delete_default_title(slide)
        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)
        # 横棒
        self.add_full_height_image(factory,slide)

        return slide

    def render_quote(self, factory, data):
        return slides_quote.render_quote_default(factory, data)

    def render_hero(self, factory, data):
        return slides_hero.render_hero_default(factory, data)
    
    def render_features(self, factory, data):
        slide = slides_features.render_features_default(factory, data)
        self.delete_default_title(slide)
        
        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)
        # 横棒
        self.add_full_height_image(factory,slide)

        return slide
       
    def render_closing(self, factory, data):
        slide = slides_closing.render_closing_default(factory, data)
        self.delete_default_title(slide)
        # タイトル
        title = data.get("title")
        self.top_title(slide, factory, title)
        # 横棒
        self.add_full_height_image(factory,slide)
