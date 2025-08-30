from typing import Any, Dict

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# 表形式
def render_table_default(factory, data: Dict[str, Any]):
    s = factory._new_slide(data)
    slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height

    factory._add_slide_title(s, data.get("title", "表"))

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    n_rows, n_cols = len(rows) + 1, len(headers)

    top = Pt(100)
    left = Pt(40)
    width = int(slide_w - Pt(80))
    height = int(slide_h * 0.55)

    table_shape = s.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 列幅（整数化必須）
    col_width = int(width / n_cols)
    for col in table.columns:
        col.width = col_width

    # 行高さ（整数化必須）
    row_height = int(height / n_rows)
    for row in table.rows:
        row.height = row_height

    # ヘッダー
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "BIZ UDゴシック"
        run.font.color.rgb = factory.colors["background"]
        cell.fill.solid()
        cell.fill.fore_color.rgb = factory.colors["primary"]
        p.alignment = PP_ALIGN.CENTER

    # データ
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(16)
            run.font.name = "BIZ UDゴシック"
            run.font.color.rgb = factory.colors["text"]
            p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = factory.colors["surface"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = factory.colors["background"]

    # bodyText（高さが溢れないように制限）
    body_text = data.get("bodyText")
    if body_text:
        b_top = min(top + height + Pt(20), slide_h - Pt(100))
        b_left = Pt(40)
        b_width = int(slide_w - Pt(80))
        b_height = int(slide_h - b_top - Pt(40))

        box = s.shapes.add_textbox(b_left, b_top, b_width, b_height)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = body_text
        run.font.size = Pt(16)
        run.font.name = "BIZ UDゴシック"
        run.font.color.rgb = factory.colors["text"]

    return s
