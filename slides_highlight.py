import math
from pptx.util import Pt
from typing import Any, Dict
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def render_highlight_default(factory, data: Dict[str, Any]):
    """
    type: "highlight"
    title: スライドタイトル
    keyword: 強調するキーワードや公式
    description: 下に配置する解説文
    """
    s = factory._new_slide(data)

    factory._add_slide_title(s, data["title"])

    slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height
    keyword = data.get("keyword", "")
    description = data.get("description", "")

    # ---------------- キーワードボックス ----------------
    font_size = 44
    line_height = font_size * 1.4
    max_chars_per_line = 14   # 1行あたりの想定文字数（日本語ベース）
    n_lines = math.ceil(len(keyword) / max_chars_per_line)

    box_w = slide_w * 0.7      # 幅は広め固定
    box_h = Pt(line_height * n_lines + 60)  # 行数に応じて高さ調整
    left = (slide_w - box_w) / 2
    top = slide_h * 0.35

    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = factory.colors["surface"]
    shape.line.color.rgb = factory.colors["accent"]

    tf = shape.text_frame
    tf.text = keyword
    p = tf.paragraphs[0]
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.name = "BIZ UDPゴシック"
    run.font.bold = True
    run.font.color.rgb = factory.colors["primary"]
    p.alignment = PP_ALIGN.CENTER

    # ---------------- 解説文 ----------------
    if description:
        desc_top = top + box_h + Pt(30)
        tbox = s.shapes.add_textbox(
            Pt(60), desc_top, slide_w - Pt(120), Pt(120)
        )
        tf_desc = tbox.text_frame
        tf_desc.word_wrap = True
        tf_desc.text = description
        p2 = tf_desc.paragraphs[0]
        run2 = p2.runs[0]
        run2.font.size = Pt(20)
        run2.font.name = "BIZ UDPゴシック"
        run2.font.color.rgb = factory.colors["text"]
        p2.alignment = PP_ALIGN.CENTER    
