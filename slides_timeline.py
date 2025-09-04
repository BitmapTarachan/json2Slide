from typing import Any, Dict

from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement


def render_timeline_default(factory, data: Dict[str, Any]):
    """タイムライン"""
    s = factory._new_slide(data)
    factory._add_slide_title(s, data.get("title",""))

    slide_height = factory.prs.slide_height

    # バー位置をスライド中央に配置
    bar_height = Cm(0.25)
    bar_top = (slide_height - bar_height) / 2   # 上下中央
    bar_left = Cm(4)                            # ← 左右を少し余白大きめに
    bar_width = factory.prs.slide_width - Cm(8)    # ← margin=4cmずつ確保

    # ベースバー
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_width, bar_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = factory.colors["ghost"]
    bar.line.fill.background()

    milestones = data.get("milestones", [])
    if not milestones:
        return s  # データがなければラインだけで終了

    step_x = bar_width / (len(milestones)-1 if len(milestones) > 1 else 1)

    for i, m in enumerate(milestones):
        title = str(m.get("label", f"Step {i+1}"))
        date = str(m.get("date", ""))

        x = bar_left + i * step_x

        # マーカー
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x-Cm(0.2), bar_top-Cm(0.2), Cm(0.4), Cm(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = factory.colors["accent"]
        circle.line.fill.background()

        # タイトル（2行分確保・下揃え）
        tbox = s.shapes.add_textbox(x-Cm(2), bar_top-Cm(2), Cm(4), Cm(1.5))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = tf.paragraphs[0]
        factory._style_text(p, title, factory.fonts["sizes"]["body"], bold=True, color=factory.colors["text"], align=PP_ALIGN.CENTER)

        # 日付
        dbox = s.shapes.add_textbox(x-Cm(2), bar_top+Cm(0.5), Cm(4), Cm(1))
        dp = dbox.text_frame.paragraphs[0]
        factory._style_text(dp, date, factory.fonts["sizes"]["caption"], color=factory.colors["subtext"], align=PP_ALIGN.CENTER)
    
    return s
