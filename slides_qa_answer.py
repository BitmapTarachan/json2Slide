from typing import Any, Dict

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

# --- Q&A: Answer ---
def render_qa_answer_default(factory, data: Dict[str, Any]):
    s = factory._new_slide(data)
    slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height

    # ゴースト "Ａ"（全角、大きく左上）
    a_size = int(min(slide_w, slide_h) * 0.5)
    abox = s.shapes.add_textbox(Pt(0), Pt(0), a_size, a_size)
    tf_a = abox.text_frame
    tf_a.word_wrap = False
    tf_a.vertical_anchor = MSO_ANCHOR.TOP

    ap = tf_a.paragraphs[0]
    run = ap.add_run()
    run.text = "Ａ"
    run.font.size = Pt(200)   # ゴーストA専用サイズ
    run.font.bold = True
    run.font.color.rgb = factory.colors["ghost"]
    ap.alignment = PP_ALIGN.LEFT

    # 答え（中央に一言）
    ans_box = s.shapes.add_textbox(Pt(100), Pt(200), slide_w - Pt(200), Pt(100))
    tf_ans = ans_box.text_frame
    ans_p = tf_ans.paragraphs[0]
    factory._style_text(
        ans_p,
        "答え : " + data.get("answer", ""),
        factory.fonts["sizes"]["contentTitle"],
        bold=True,
        color=factory.colors["primary"],
        align=PP_ALIGN.CENTER
    )

    # 解説（中央寄せ）
    exp_box = s.shapes.add_textbox(Pt(100), slide_h/3, slide_w - Pt(200), slide_h/2)
    tf_exp = exp_box.text_frame
    tf_exp.word_wrap = True
    tf_exp.vertical_anchor = MSO_ANCHOR.MIDDLE
    exp_p = tf_exp.paragraphs[0]
    factory._style_text(
        exp_p,
        data.get("explanation", ""),
        factory.fonts["sizes"]["body"],
        color=factory.colors["text"],
        align=PP_ALIGN.CENTER
    )

    return s
