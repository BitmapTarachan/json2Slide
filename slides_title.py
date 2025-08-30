#slides_title.py

from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE

def render_title_default(factory, data):
    s = factory._new_slide(data)
    slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height

    # 右側に画像を配置（オプション）
    if "image" in data and data["image"]:
        stream, im = factory._load_image(data["image"])
        if stream and im:
            iw, ih = im.size
            slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height

            # 縦に合わせる（はみ出さないように調整）
            scale = slide_h / ih
            new_w, new_h = iw * scale, ih * scale

            left = int(slide_w * 3/5)
            top = 0
            if left + new_w < slide_w:  # 幅足りないなら右寄せ
                left = slide_w - new_w

            s.shapes.add_picture(stream, left, top, width=int(new_w), height=int(new_h))

    # 左側の縦線（Ghostカラー）
    line_left = Pt(50)  # 余白を少し空ける
    s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        line_left, Pt(90),  # X, Y
        Pt(2), slide_h - Pt(200)  
    ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = factory.colors["ghost"]
    s.shapes[-1].line.fill.background()  # 枠線なし
    s.shapes[-1].shadow.inherit = False  # 影無し

    # 教科名
    subj_rect = factory.layout.get_rect("titleSlide.subject")
    subj_box = s.shapes.add_textbox(subj_rect["left"], subj_rect["top"], subj_rect["width"], subj_rect["height"])
    subj_p = subj_box.text_frame.paragraphs[0]
    factory._style_text(
        subj_p,
        data.get("subject", ""),
        factory.fonts["sizes"]["contentTitle"],
        color=factory.colors["subtext"]
    )

    # タイトル
    rect = factory.layout.get_rect("titleSlide.title")
    box = s.shapes.add_textbox(rect["left"], rect["top"], rect["width"], rect["height"])
    p = box.text_frame.paragraphs[0]
    factory._style_text(
        p,
        data.get("title", ""),
        factory.fonts["sizes"]["title"],
        bold=True,
        color=factory.colors["primary"]
    )

    # 講師名（固定）
    l_rect = factory.layout.get_rect("titleSlide.lecturer")
    lbox = s.shapes.add_textbox(l_rect["left"], l_rect["top"], l_rect["width"], l_rect["height"])
    lp = lbox.text_frame.paragraphs[0]
    factory._style_text(
        lp,
        "講師名：〇〇　〇〇",
        factory.fonts["sizes"]["body"],
        color=factory.colors["text"]
    )

    # 日付
    d_rect = factory.layout.get_rect("titleSlide.date")
    dbox = s.shapes.add_textbox(d_rect["left"], d_rect["top"], d_rect["width"], d_rect["height"])
    dp = dbox.text_frame.paragraphs[0]
    factory._style_text(
        dp,
        data.get("date", ""),
        factory.fonts["sizes"]["body"],
        color=factory.colors["subtext"]
    )

    return s
