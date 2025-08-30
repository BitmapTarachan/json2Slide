# slides_quote.py
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE

def render_quote_default(factory, data):
    s = factory._new_slide(data,False)
    slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height
    side_w = slide_w / 3

    # --- 左半分（画像 or primaryカラー） ---
    bg_url = data.get("image")
    if bg_url:
        stream, im = factory._load_image(bg_url)
        if stream:
            iw, ih = im.size
            slide_h = factory.prs.slide_height

            # 縦をスライドにフィット
            scale = slide_h / ih
            new_w = iw * scale
            new_h = ih * scale
                
            s.shapes.add_picture(stream, 0, 0, width=new_w, height=new_h)

        # 半透明オーバーレイ
        overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, side_w, slide_h)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = factory.colors["ghost"]
        factory._set_shape_transparency(overlay, 40000)  # 40%透過
        overlay.line.fill.background()
        overlay.shadow.inherit = False
    else:
        # 画像がなければprimaryカラーで塗りつぶし
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, side_w, slide_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = factory.colors["ghost"]
        rect.line.fill.background()
        rect.shadow.inherit = False

    # --- 右半分（背景はスライド標準色） ---
    rect_right = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, side_w, 0, side_w * 2, slide_h)
    rect_right.fill.solid()
    rect_right.fill.fore_color.rgb = factory.colors["background"]
    rect_right.line.fill.background()
    rect_right.shadow.inherit = False

    # --- 引用符アイコン ---
    quote_box = s.shapes.add_textbox(side_w + Pt(40), Pt(40), side_w - Pt(80), Pt(80))
    qf = quote_box.text_frame
    qp = qf.paragraphs[0]
    qp.text = "“"   # フォント依存でシャープな形を狙う
    run = qp.runs[0]
    run.font.name = "Arial"   
    run.font.size = Pt(150)
    run.font.bold = True
    run.font.color.rgb = factory.colors["ghost"]
    qp.alignment = PP_ALIGN.LEFT

    # --- 引用文 ---
    qbox = s.shapes.add_textbox(side_w + Pt(80), slide_h*0.25, side_w * 2 - Pt(160), slide_h*0.4)
    tf = qbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    factory._style_text(
        p,
        data.get("quote", ""),
        factory.fonts["sizes"]["sectionTitle"],
        bold=True,
        color=factory.colors["text"],
        align=PP_ALIGN.LEFT
    )

    # --- 引用元 ---
    author = data.get("author", "")
    if author:
        abox = s.shapes.add_textbox(side_w + Pt(80), slide_h*0.75, side_w * 2 - Pt(80), Pt(80))
        atf = abox.text_frame
        ap = atf.paragraphs[0]
        factory._style_text(
            ap,
            author,
            factory.fonts["sizes"]["subhead"],
            color=factory.colors["subtext"],
            align=PP_ALIGN.LEFT
        )

    return s
