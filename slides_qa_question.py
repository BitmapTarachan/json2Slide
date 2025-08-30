from typing import Any, Dict

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

# --- Q&A: Question ---
def render_qa_question_defaults(factory, data: Dict[str, Any]):
    s = factory._new_slide(data)
    slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height

    # ゴースト "Ｑ"（全角、大きく左上）
    q_size = int(min(slide_w, slide_h) * 0.5)  # スライドの半分くらいの大きさ
    qbox = s.shapes.add_textbox(Pt(0), Pt(0), q_size, q_size)
    tf_q = qbox.text_frame
    tf_q.word_wrap = False
    tf_q.vertical_anchor = MSO_ANCHOR.TOP

    qp = tf_q.paragraphs[0]
    run = qp.add_run()
    run.text = "Ｑ"
    run.font.size = Pt(200)   # ゴーストQ専用サイズ（必要に応じて調整）
    run.font.bold = True
    run.font.color.rgb = factory.colors["ghost"]
    qp.alignment = PP_ALIGN.LEFT

    # 質問文（中央揃え）
    qtext_box = s.shapes.add_textbox(Pt(100), slide_h/3, slide_w - Pt(200), slide_h/3)
    tf = qtext_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    factory._style_text(
        p,
        data.get("question", ""),
        factory.fonts["sizes"]["sectionTitle"],
        bold=True,
        color=factory.colors["text"],
        align=PP_ALIGN.CENTER
    )

    return s
