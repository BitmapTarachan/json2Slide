#slides_hero.py

from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def render_hero_default(factory, data):
    s = factory._new_slide(data)

    # スライドサイズ
    slide_w, slide_h = factory.prs.slide_width, factory.prs.slide_height

    # 背景画像（あれば適用）
    bg_url = data.get("background-image")
    if bg_url:
        stream, _ = factory._load_image(bg_url)
        if stream:
            s.shapes.add_picture(stream, 0, 0, width=slide_w, height=slide_h)
    else:
        # 背景色だけ
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = factory.colors["background"]

    # --- 半透明オーバーレイ ---
    overlay = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 黒
    overlay.fill.fore_color.transparency = 0.5              # 30%透過
    overlay.line.fill.background()                   # 枠線なし
    factory._set_shape_transparency(overlay, 40000)  # 40%ほど透過

    # --- タイトル（中央配置） ---
    tbox = s.shapes.add_textbox(0, slide_h*0.35, slide_w, slide_h*0.2)
    tf = tbox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tp = tf.paragraphs[0]
    factory._style_text(
        tp,
        data.get("title", ""),
        factory.fonts["sizes"]["title"],
        bold=True,
        color=RGBColor(255, 255, 255),
        align=PP_ALIGN.CENTER
    )

    # --- 仕切り線 ---
    line_top = slide_h * 0.51
    line_left = slide_w * 0.1
    line_width = slide_w * 0.8
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()

    # --- サブタイトル（タイトルのすぐ下） ---
    subtitle = data.get("subtitle")
    if subtitle:
        stbox = s.shapes.add_textbox(0, line_top + Pt(10), slide_w, slide_h*0.15)
        stf = stbox.text_frame
        stf.vertical_anchor = MSO_ANCHOR.TOP
        sp = stf.paragraphs[0]
        factory._style_text(
            sp,
            subtitle,
            factory.fonts["sizes"]["subhead"],
            color=RGBColor(230, 230, 230),
            align=PP_ALIGN.CENTER
        )

    return s
