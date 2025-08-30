from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt

def render_image3_default(self, slide, images, font_size):
    slide_w, slide_h = self.prs.slide_width, self.prs.slide_height
    margin = Pt(40)
    spacing = Pt(30)
    max_width_ratio = 0.8  # 横幅全体の80%に収める

    # 各画像の幅（等分）
    target_w = (slide_w * max_width_ratio - 2 * spacing) / 3

    scaled_sizes = []
    streams = []
    for img in images:
        stream, im = self._load_image(img["url"])
        if stream :
            iw, ih = im.size
            scale = target_w / iw
            new_w, new_h = iw * scale, ih * scale
            scaled_sizes.append((int(new_w), int(new_h)))
            streams.append(stream)

    # 横方向の開始位置（中央寄せ）
    total_w = sum(w for w, h in scaled_sizes) + 2 * spacing
    left_start = (slide_w - total_w) / 2

    # 最大高さを揃える（縦位置は上を揃える）
    img_max_h = max(h for w, h in scaled_sizes)
    top_img = (slide_h / 2) - img_max_h / 2 - Pt(20)

    # キャプションのY座標（全画像共通で揃える）
    cap_top = top_img + img_max_h + Pt(10)

    shapes = []
    x = left_start
    for (w, h), stream, img in zip(scaled_sizes, streams, images):
        # 画像（上を揃える。小さい画像は下に余白）
        top = top_img + (img_max_h - h)
        if stream:
            slide.shapes.add_picture(stream, x, top, width=w, height=h)

        # キャプション（Y位置を固定）
        cap_box = slide.shapes.add_textbox(x, cap_top, w, Pt(40))
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        self._style_text(p, img.get("caption", ""), font_size, self.colors["text"])
        p.alignment = PP_ALIGN.LEFT
        tf.vertical_anchor = MSO_ANCHOR.TOP
        shapes.append(cap_box)

        x += w + spacing

    return shapes
