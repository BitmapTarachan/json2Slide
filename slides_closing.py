from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Cm

def render_closing_default(factory, data):
    s = factory._new_slide(data)

    factory._add_slide_title(s, data["title"])

    # 本文エリアの基準
    body_left = Cm(2.0)
    body_top = Pt(100) + Cm(1.0)
    body_width = factory.prs.slide_width * 0.55  # 左側にテキスト
    body_height = factory.prs.slide_height - body_top - Cm(2.0)

    # headline（強調文）
    headline = data.get("headline", "")
    if headline:
        hbox = s.shapes.add_textbox(body_left, body_top, body_width, Cm(3))
        tf_h = hbox.text_frame
        tf_h.clear()
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        factory._style_text(
            p_h,
            headline,
            Pt(28),
            bold=True,
            color=factory.colors["text"]
        )

    # bodyText（補足説明）
    body_text = data.get("bodyText", "")
    if body_text:
        bbox = s.shapes.add_textbox(body_left, body_top + Cm(3.5), body_width, body_height - Cm(3.5))
        tf_b = bbox.text_frame
        tf_b.clear()
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        factory._style_text(
            p_b,
            body_text,
            Pt(18),
            color=factory.colors["text"]
        )
    text_bottom = bbox.top +bbox.height

    # 右側の画像
    if "image" in data and data["image"]:
        stream, im = factory._load_image(data["image"])
        if stream and im:
            iw, ih = im.size

            # テキスト領域の上下に合わせる
            img_top = body_top                     # headline の開始位置
            img_bottom = text_bottom
            avail_h = img_bottom - img_top

            # 縦に合わせてスケーリング
            scale = avail_h / ih
            new_w, new_h = iw * scale, ih * scale

            left = body_left + body_width + Cm(1.0)
            s.shapes.add_picture(stream, left, img_top, width=int(new_w), height=int(new_h))

    return s
