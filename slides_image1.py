from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt

# --- 1枚: 画像 + 右にキャプション ---
def render_image1_default(factory, slide, images, font_size):
    img = images[0]
    caption = img.get("caption", "")

    # スライドサイズ
    slide_width = factory.prs.slide_width
    slide_height = factory.prs.slide_height
    margin = Pt(20)

    # factory._load_image で画像読み込み（BytesIO, PIL.Image）
    stream, im = factory._load_image(img["url"])
    if stream :
        iw, ih = im.size
        aspect = iw / ih

        # 横長か縦長かでリサイズ基準を切替
        if aspect >= 1:  # 横長 → 横幅を中央まで広げる
            max_width = slide_width / 2 - 2 * margin
            width = max_width
            height = width / aspect
            left = margin
            top = (slide_height - height) / 2
        else:  # 縦長 → 上下いっぱいまで
            max_height = slide_height - 2 * margin
            height = max_height
            width = height * aspect
            left = margin
            top = (slide_height - height) / 2

        # 画像挿入
        slide.shapes.add_picture(stream, left, top, width=width, height=height)

    # キャプションテキスト
    cap_left = slide_width / 2 + margin
    cap_width = slide_width / 2 - 2 * margin
    cap_height = slide_height - 2 * margin
    cap_top = margin

    txBox = slide.shapes.add_textbox(cap_left, cap_top, cap_width, cap_height)
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.alignment = PP_ALIGN.LEFT
    factory._style_text(p, caption, font_size, factory.colors["text"])

